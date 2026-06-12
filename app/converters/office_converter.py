from __future__ import annotations

from pathlib import Path
import shutil
import subprocess
import tempfile
from shutil import which

from app.converters.pdf_converter import POPPLER_ERROR_MESSAGE, pdf_to_png
from app.services.file_service import (
    OUTPUT_DIR,
    ensure_directory,
    is_supported_pptx,
    unique_output_path,
)


LIBREOFFICE_ERROR_MESSAGE = (
    "PPTX to PDF requires LibreOffice. On macOS, install it with: "
    "brew install --cask libreoffice"
)
PPTX_SLIDE_IMAGE_EXPORT_ERROR_MESSAGE = (
    "PPTX slide image export requires LibreOffice and Poppler. On macOS, "
    "install them with: brew install --cask libreoffice and brew install poppler"
)
MACOS_SOFFICE_PATH = Path("/Applications/LibreOffice.app/Contents/MacOS/soffice")
PPTX_DOCX_MODE_TEXT_OUTLINE = "text_outline"
PPTX_DOCX_MODE_SLIDE_IMAGES = "slide_images"
PPTX_DOCX_MODE_SLIDE_IMAGES_WITH_TEXT = "slide_images_with_text"
PPTX_DOCX_MODES = {
    PPTX_DOCX_MODE_TEXT_OUTLINE,
    PPTX_DOCX_MODE_SLIDE_IMAGES,
    PPTX_DOCX_MODE_SLIDE_IMAGES_WITH_TEXT,
}
PPTX_DOCX_MODE_ALIASES = {
    "text outline": PPTX_DOCX_MODE_TEXT_OUTLINE,
    "slide images": PPTX_DOCX_MODE_SLIDE_IMAGES,
    "slide images + extracted text": PPTX_DOCX_MODE_SLIDE_IMAGES_WITH_TEXT,
    "slide images with text": PPTX_DOCX_MODE_SLIDE_IMAGES_WITH_TEXT,
}


def soffice_executable() -> str | None:
    executable = which("soffice")
    if executable:
        return executable

    if MACOS_SOFFICE_PATH.exists():
        return str(MACOS_SOFFICE_PATH)

    return None


def ensure_soffice_available() -> str:
    executable = soffice_executable()
    if not executable:
        raise RuntimeError(LIBREOFFICE_ERROR_MESSAGE)

    return executable


def pptx_to_pdf(
    input_path: str | Path,
    output_dir: str | Path = OUTPUT_DIR,
) -> Path:
    source = Path(input_path)
    if not is_supported_pptx(source):
        raise ValueError(f"Unsupported PowerPoint file: {source.name}")

    executable = ensure_soffice_available()
    output_directory = ensure_directory(output_dir)

    with tempfile.TemporaryDirectory(dir=output_directory) as temporary_directory:
        try:
            subprocess.run(
                [
                    executable,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    temporary_directory,
                    str(source),
                ],
                check=True,
                capture_output=True,
                text=True,
            )
        except subprocess.CalledProcessError as exc:
            details = (exc.stderr or exc.stdout or "").strip()
            if exc.returncode == 127 or "No such file or directory" in details:
                raise RuntimeError(LIBREOFFICE_ERROR_MESSAGE) from exc

            message = "PPTX to PDF conversion failed."
            if details:
                message = f"{message} {details}"
            raise RuntimeError(message) from exc

        converted_path = Path(temporary_directory) / f"{source.stem}.pdf"
        if not converted_path.exists():
            raise RuntimeError("PPTX to PDF conversion failed: no PDF was generated.")

        output_path = unique_output_path(source, "pdf", output_directory)
        shutil.move(str(converted_path), output_path)

    return output_path


def normalize_pptx_docx_mode(mode: str) -> str:
    label = " ".join(mode.lower().strip().replace("-", " ").split())
    if label in PPTX_DOCX_MODE_ALIASES:
        return PPTX_DOCX_MODE_ALIASES[label]

    normalized = mode.lower().strip().replace("-", "_").replace(" ", "_")
    if normalized not in PPTX_DOCX_MODES:
        raise ValueError(f"Unsupported PPTX to DOCX mode: {mode}")

    return normalized


def export_pptx_slide_images(
    input_path: str | Path,
    output_dir: str | Path = OUTPUT_DIR,
    dpi: int = 200,
) -> list[Path]:
    source = Path(input_path)
    if not is_supported_pptx(source):
        raise ValueError(f"Unsupported PowerPoint file: {source.name}")

    output_directory = ensure_directory(output_dir)
    with tempfile.TemporaryDirectory(dir=output_directory) as temporary_directory:
        temporary_path = Path(temporary_directory)
        pdf_path = pptx_to_pdf(source, temporary_path)
        try:
            image_paths = pdf_to_png(pdf_path, temporary_path, dpi=dpi)
        except RuntimeError as exc:
            if str(exc) == POPPLER_ERROR_MESSAGE:
                raise RuntimeError(PPTX_SLIDE_IMAGE_EXPORT_ERROR_MESSAGE) from exc
            raise

        if not image_paths:
            raise RuntimeError(
                "PPTX slide image export failed: no slide images were generated."
            )

        copied_paths: list[Path] = []
        for index, image_path in enumerate(image_paths, start=1):
            output_path = unique_output_path(
                f"{source.stem}-slide-{index}.png",
                "png",
                output_directory,
            )
            shutil.copy2(image_path, output_path)
            copied_paths.append(output_path)

    return copied_paths


def clean_text_lines(text: str) -> list[str]:
    return [line.strip() for line in text.splitlines() if line.strip()]


def text_lines_from_shape(shape) -> list[str]:
    if not shape or not getattr(shape, "has_text_frame", False):
        return []

    text_frame = getattr(shape, "text_frame", None)
    return clean_text_lines(getattr(text_frame, "text", ""))


def speaker_note_lines(slide) -> list[str]:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return []

    text_frame = getattr(notes_slide, "notes_text_frame", None)
    return clean_text_lines(getattr(text_frame, "text", ""))


def slide_text_sections(slide) -> tuple[str, list[str], list[str]]:
    shapes = getattr(slide, "shapes", [])
    title_shape = getattr(shapes, "title", None)
    title_lines = text_lines_from_shape(title_shape)
    title = title_lines[0] if title_lines else ""
    content_lines: list[str] = []

    for shape in shapes:
        if shape is title_shape:
            continue
        content_lines.extend(text_lines_from_shape(shape))

    return title, content_lines, speaker_note_lines(slide)


def add_slide_text_to_document(document, slide, include_title: bool = True) -> None:
    title, content_lines, notes_lines = slide_text_sections(slide)
    if include_title:
        document.add_paragraph(f"Title: {title}")

    document.add_paragraph("Content:")
    if content_lines:
        for line in content_lines:
            document.add_paragraph(line, style="List Bullet")
    else:
        document.add_paragraph("(No slide body text found.)")

    if notes_lines:
        document.add_paragraph("Speaker Notes:")
        for line in notes_lines:
            document.add_paragraph(line, style="List Bullet")


def load_presentation_slides(input_path: str | Path):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for PPTX text extraction."
        ) from exc

    return Presentation(str(input_path)).slides


def new_docx_document():
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for DOCX output.") from exc

    return Document()


def pptx_to_docx_text_outline(
    input_path: str | Path,
    output_dir: str | Path = OUTPUT_DIR,
) -> Path:
    source = Path(input_path)
    if not is_supported_pptx(source):
        raise ValueError(f"Unsupported PowerPoint file: {source.name}")

    output_path = unique_output_path(source, "docx", output_dir)
    slides = load_presentation_slides(source)
    document = new_docx_document()

    for index, slide in enumerate(slides, start=1):
        if index > 1:
            document.add_page_break()

        document.add_heading(f"Slide {index}", level=1)
        add_slide_text_to_document(document, slide)

    document.save(output_path)
    return output_path


def add_slide_image_to_document(document, image_path: Path) -> None:
    section = document.sections[-1]
    max_width = section.page_width - section.left_margin - section.right_margin
    document.add_picture(str(image_path), width=max_width)


def pptx_to_docx_with_slide_images(
    input_path: str | Path,
    output_dir: str | Path = OUTPUT_DIR,
    include_extracted_text: bool = False,
) -> Path:
    source = Path(input_path)
    if not is_supported_pptx(source):
        raise ValueError(f"Unsupported PowerPoint file: {source.name}")

    output_path = unique_output_path(source, "docx", output_dir)
    document = new_docx_document()
    slides = load_presentation_slides(source) if include_extracted_text else []

    with tempfile.TemporaryDirectory(dir=ensure_directory(output_dir)) as temporary_dir:
        slide_images = export_pptx_slide_images(source, temporary_dir)

        for index, image_path in enumerate(slide_images, start=1):
            if index > 1:
                document.add_page_break()

            document.add_heading(f"Slide {index}", level=1)
            add_slide_image_to_document(document, image_path)

            if include_extracted_text:
                document.add_heading("Extracted Text", level=2)
                if index <= len(slides):
                    add_slide_text_to_document(document, slides[index - 1])
                else:
                    document.add_paragraph("(No slide text found.)")

        document.save(output_path)

    return output_path


def pptx_to_docx(
    input_path: str | Path,
    output_dir: str | Path = OUTPUT_DIR,
    mode: str = PPTX_DOCX_MODE_TEXT_OUTLINE,
) -> Path:
    normalized_mode = normalize_pptx_docx_mode(mode)

    if normalized_mode == PPTX_DOCX_MODE_TEXT_OUTLINE:
        return pptx_to_docx_text_outline(input_path, output_dir)

    return pptx_to_docx_with_slide_images(
        input_path,
        output_dir,
        include_extracted_text=(
            normalized_mode == PPTX_DOCX_MODE_SLIDE_IMAGES_WITH_TEXT
        ),
    )
